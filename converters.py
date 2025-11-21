import os
import re
import requests
import traceback
import logging
import tempfile
import textwrap
from pathlib import Path

import markdown  # For markdown to HTML conversion

logger = logging.getLogger(__name__)
  
def convert_html_to_pdf(file_path, output_path):
    """
    Convert HTML file to PDF using Gotenberg when available, otherwise fall back to a
    pure-Python implementation backed by ReportLab.

    Args:
        file_path (str): Path to the HTML file
        output_path (str): Path where the PDF will be saved

    Returns:
        bool: True if conversion successful, False otherwise
    """
    GOTENBERG_BASE_URL = os.getenv("PDF_CONVERT_URL", "http://localhost:7777")
    api_url = f"{GOTENBERG_BASE_URL}/forms/chromium/convert/html"
    success = False

    try:
        logger.info(f"Converting HTML at {file_path} to PDF, file exists = {os.path.exists(file_path)}")
        assert os.path.exists(file_path), f"HTML file not found: {file_path}"

        if not file_path.lower().endswith(('.html', '.htm')):
            raise ValueError(f"Expected HTML file, got: {file_path}")

        with open(file_path, 'rb') as f:
            files = {'files': ('index.html', f, 'text/html')}
            payload = {
                'paperWidth': '8.27',
                'paperHeight': '11.7',
                'marginTop': '0.39',
                'marginBottom': '0.39',
                'marginLeft': '0.39',
                'marginRight': '0.39',
                'preferCSSPageSize': 'false',
                'printBackground': 'true',
                'landscape': 'false',
                'scale': '1.0'
            }

            logger.info(f"Sending HTML conversion request to: {api_url}")
            response = requests.post(
                api_url,
                files=files,
                data=payload,
                timeout=60
            )

            if response.status_code == 200:
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, 'wb') as out_file:
                    out_file.write(response.content)

                logger.info(f"✅ HTML to PDF conversion successful: {output_path}")
                success = True
            else:
                logger.error(f"❌ HTML conversion failed with status {response.status_code}: {response.text}")

    except Exception as e:
        exc = traceback.format_exc()
        logger.error(f"❌ Exception converting HTML at {file_path} to PDF: {e}\n{exc}")

    if success:
        return True

    logger.warning("Falling back to pure-Python HTML to PDF conversion for %s", file_path)
    return _fallback_convert_html_file(file_path, output_path)


def convert_doc_to_pdf(file_path, output_path):
    """
    Convert Office documents to PDF via Gotenberg when accessible, or fall back to a
    simplified text-based conversion implemented in Python.

    Args:
        file_path (str): Path to the document
        output_path (str): Path where the PDF will be saved

    Returns:
        bool: True if conversion successful, False otherwise
    """
    GOTENBERG_BASE_URL = os.getenv("PDF_CONVERT_URL", "http://localhost:7777")
    api_url = f"{GOTENBERG_BASE_URL}/forms/libreoffice/convert"
    supported_extensions = ('.docx', '.doc', '.odt', '.rtf', '.xlsx', '.xls', '.pptx', '.ppt')
    success = False

    try:
        logger.info(f"Converting DOCX at {file_path} to PDF, file exists = {os.path.exists(file_path)}")
        assert os.path.exists(file_path), f"DOCX file not found: {file_path}"

        if not file_path.lower().endswith(supported_extensions):
            raise ValueError(f"Unsupported file type. Expected: {supported_extensions}, got: {file_path}")

        mime_types = {
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.odt': 'application/vnd.oasis.opendocument.text',
            '.rtf': 'application/rtf',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.xls': 'application/vnd.ms-excel',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint'
        }

        file_ext = os.path.splitext(file_path)[1].lower()
        mime_type = mime_types.get(file_ext, 'application/octet-stream')

        with open(file_path, 'rb') as f:
            files = {'files': (os.path.basename(file_path), f, mime_type)}
            payload = {
                'pdfFormat': 'PDF/A-1a',
                'merge': 'false',
                'landscape': 'false',
                'nativePageRanges': '',
                'exportFormFields': 'true',
                'allowDuplicateFieldNames': 'false',
                'exportBookmarks': 'true',
                'exportBookmarksToPdfDestination': 'false',
                'exportPlaceholders': 'false',
                'exportNotes': 'false',
                'exportNotesPages': 'false',
                'exportOnlyNotesPages': 'false',
                'exportNotesInMargin': 'false',
                'convertOooTargetToPdfTarget': 'false',
                'exportLinksRelativeFsys': 'false',
                'exportHiddenSlides': 'false',
                'skipEmptyPages': 'true',
                'addOriginalDocumentAsStream': 'false'
            }

            logger.info(f"Sending DOCX conversion request to: {api_url}")
            response = requests.post(
                api_url,
                files=files,
                data=payload,
                timeout=120
            )

            if response.status_code == 200:
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, 'wb') as out_file:
                    out_file.write(response.content)

                logger.info(f"✅ DOCX to PDF conversion successful: {output_path}")
                success = True
            else:
                logger.error(f"❌ DOCX conversion failed with status {response.status_code}: {response.text}")

    except Exception as e:
        exc = traceback.format_exc()
        logger.error(f"❌ Exception converting DOCX at {file_path} to PDF: {e}\n{exc}")

    if success:
        return True

    logger.warning("Falling back to pure-Python document conversion for %s", file_path)
    return _fallback_convert_doc_to_pdf(file_path, output_path)


def convert_html_string_to_pdf(html_content, output_path, title="Document"):
    """
    Convert HTML string content to PDF, preferring Gotenberg but gracefully falling back
    to a ReportLab-based implementation when the service is unavailable.
    """
    GOTENBERG_BASE_URL = os.getenv("PDF_CONVERT_URL", "http://localhost:7777")
    api_url = f"{GOTENBERG_BASE_URL}/forms/chromium/convert/html"
    success = False

    try:
        logger.info(f"Converting HTML string to PDF at {output_path}")
        files = {'files': ('index.html', html_content.encode('utf-8'), 'text/html')}
        payload = {
            'paperWidth': '8.27',
            'paperHeight': '11.7',
            'marginTop': '0.39',
            'marginBottom': '0.39',
            'marginLeft': '0.39',
            'marginRight': '0.39',
            'preferCSSPageSize': 'false',
            'printBackground': 'true',
            'landscape': 'false',
            'scale': '1.0'
        }

        logger.info(f"Sending HTML string conversion request to: {api_url}")
        response = requests.post(
            api_url,
            files=files,
            data=payload,
            timeout=60
        )

        if response.status_code == 200:
            os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
            with open(output_path, 'wb') as out_file:
                out_file.write(response.content)

            logger.info(f"✅ HTML string to PDF conversion successful: {output_path}")
            success = True
        else:
            logger.error(f"❌ HTML string conversion failed with status {response.status_code}: {response.text}")

    except Exception as e:
        exc = traceback.format_exc()
        logger.error(f"❌ Exception converting HTML string to PDF: {e}\n{exc}")

    if success:
        return True

    logger.warning("Falling back to pure-Python HTML string conversion for %s", output_path)
    return _fallback_convert_html_string(html_content, output_path, title=title)


def convert_markdown_file_to_pdf(file_path, output_path):
    """
    Convert Markdown file to PDF using Gotenberg.
    
    This function reads a Markdown file, converts it to HTML using the Python markdown library,
    then uses Gotenberg's Chromium engine to convert the HTML to PDF.
    
    Args:
        file_path (str): Path to the Markdown file
        output_path (str): Path where the PDF will be saved
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    try:
        logger.info(f"Converting Markdown file at {file_path} to PDF, file exists = {os.path.exists(file_path)}")
        assert os.path.exists(file_path), f"Markdown file not found: {file_path}"
        
        # Validate file extension
        if not file_path.lower().endswith(('.md', '.markdown', '.mdown', '.mkd', '.mdwn')):
            raise ValueError(f"Expected Markdown file, got: {file_path}")
        
        # Read the markdown file
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        # Get the title from the filename
        title = os.path.splitext(os.path.basename(file_path))[0]
        
        # Use the markdown string converter
        return convert_markdown_string_to_pdf(md_content, output_path, title)
        
    except Exception as e:
        exc = traceback.format_exc()
        logger.error(f"❌ Exception converting Markdown file at {file_path} to PDF: {e}\n{exc}")
        return False


def convert_markdown_string_to_pdf(markdown_content, output_path, title="Document", css_style=None):
    """
    Convert Markdown string content to PDF by first converting to HTML then using Gotenberg.
    
    This function converts Markdown to HTML using the Python markdown library,
    then uses Gotenberg's Chromium engine to convert the HTML to PDF.
    
    Args:
        markdown_content (str): Markdown content as string
        output_path (str): Path where the PDF will be saved
        title (str): Title for the document
        css_style (str, optional): Custom CSS styles to apply to the HTML
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    try:
        md_content = markdown_content
        
        logger.info(f"Converting Markdown to PDF at {output_path}")
        
        # Convert Markdown to HTML using Python markdown library
        # Enable useful extensions for better formatting
        md = markdown.Markdown(extensions=[
            'extra',        # Includes tables, footnotes, abbreviations, etc.
            'codehilite',   # Code syntax highlighting
            'toc',          # Table of contents
            'meta',         # Metadata support
            'sane_lists',   # Better list handling
            'nl2br',        # New line to break
            'smarty',       # Smart quotes and dashes
        ])
        
        html_body = md.convert(md_content)
        
        # Default CSS for better formatting
        default_css = """
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
            line-height: 1.6;
            color: #333;
            max-width: 900px;
            margin: 0 auto;
            padding: 20px;
        }
        h1, h2, h3, h4, h5, h6 {
            margin-top: 24px;
            margin-bottom: 16px;
            font-weight: 600;
            line-height: 1.25;
        }
        h1 { font-size: 2em; border-bottom: 1px solid #eaecef; padding-bottom: 0.3em; }
        h2 { font-size: 1.5em; border-bottom: 1px solid #eaecef; padding-bottom: 0.3em; }
        h3 { font-size: 1.25em; }
        code {
            background-color: #f6f8fa;
            padding: 0.2em 0.4em;
            border-radius: 3px;
            font-family: 'Consolas', 'Monaco', 'Courier New', monospace;
        }
        pre {
            background-color: #f6f8fa;
            padding: 16px;
            overflow: auto;
            border-radius: 6px;
        }
        pre code {
            background-color: transparent;
            padding: 0;
        }
        blockquote {
            padding: 0 1em;
            color: #6a737d;
            border-left: 0.25em solid #dfe2e5;
            margin: 0;
        }
        table {
            border-collapse: collapse;
            width: 100%;
            margin: 15px 0;
        }
        table th, table td {
            padding: 6px 13px;
            border: 1px solid #dfe2e5;
        }
        table tr:nth-child(2n) {
            background-color: #f6f8fa;
        }
        img {
            max-width: 100%;
            height: auto;
        }
        a {
            color: #0366d6;
            text-decoration: none;
        }
        a:hover {
            text-decoration: underline;
        }
        """
        
        # Use custom CSS if provided, otherwise use default
        final_css = css_style if css_style else default_css
        
        # Create complete HTML document
        html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        {final_css}
    </style>
</head>
<body>
    {html_body}
</body>
</html>"""
        
        success = convert_html_string_to_pdf(html_content, output_path, title)
        if success:
            return True

        logger.warning("Markdown HTML conversion failed, falling back to text-based PDF generation.")
        return _fallback_convert_markdown_text(markdown_content, output_path, title)

    except ImportError:
        logger.error("❌ Python 'markdown' library not installed. Install it using: pip install markdown")
        if _convert_markdown_via_gotenberg(markdown_content, output_path, title):
            return True

        logger.warning("Falling back to basic Markdown text conversion for %s", output_path)
        return _fallback_convert_markdown_text(markdown_content, output_path, title)

    except Exception as e:
        exc = traceback.format_exc()
        logger.error(f"❌ Exception converting Markdown to PDF: {e}\n{exc}")
        return _fallback_convert_markdown_text(markdown_content, output_path, title)


def _fallback_convert_html_file(file_path: str, output_path: str) -> bool:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()
    except Exception as e:
        logger.error(f"❌ Failed to read HTML file for fallback conversion: {e}")
        return False

    title = Path(file_path).stem or "Document"
    return _fallback_convert_html_string(html_content, output_path, title=title)


def _fallback_convert_html_string(html_content: str, output_path: str, title: str = "Document") -> bool:
    text = _html_to_text(html_content)
    return _fallback_convert_text_to_pdf(text, output_path, title=title)


def _fallback_convert_doc_to_pdf(file_path: str, output_path: str) -> bool:
    try:
        text = _extract_text_from_document(file_path)
    except Exception as e:
        logger.error(f"❌ Fallback document conversion failed for {file_path}: {e}")
        return False

    title = Path(file_path).stem or "Document"
    return _fallback_convert_text_to_pdf(text, output_path, title=title)


def _extract_text_from_document(file_path: str) -> str:
    suffix = Path(file_path).suffix.lower()
    if suffix in (".docx", ".docm"):
        try:
            from docx import Document  # type: ignore
        except ImportError as exc:
            raise RuntimeError("python-docx is required for DOCX fallback conversion.") from exc

        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    if suffix == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text  # type: ignore
        except ImportError as exc:
            raise RuntimeError("striprtf is required for RTF fallback conversion.") from exc

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return rtf_to_text(f.read())

    if suffix in (".pptx",):
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise RuntimeError("python-pptx is required for PPTX fallback conversion.") from exc

        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    texts.append(text.strip())
        return "\n\n".join(texts)

    if suffix in (".xlsx", ".xls"):
        try:
            import pandas as pd  # type: ignore
        except ImportError as exc:
            raise RuntimeError("pandas is required for spreadsheet fallback conversion.") from exc

        df = pd.read_excel(file_path, engine="openpyxl" if suffix == ".xlsx" else None)
        return df.to_csv(index=False)

    if suffix in (".txt",):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    raise RuntimeError(f"Fallback converter does not support '{suffix}' files. Please install Gotenberg or convert the document to DOCX.")


def _fallback_convert_text_to_pdf(text: str, output_path: str, title: str = "Document") -> bool:
    if not text.strip():
        text = "No textual content was available for conversion."

    try:
        from reportlab.lib.pagesizes import LETTER  # type: ignore
        from reportlab.pdfgen import canvas  # type: ignore
    except ImportError as exc:
        logger.error("ReportLab is required for fallback PDF generation. Install it via 'pip install reportlab'.")
        return False

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    pdf_canvas = canvas.Canvas(output_path, pagesize=LETTER)
    pdf_canvas.setTitle(title)
    pdf_canvas.setFont("Helvetica", 11)

    width, height = LETTER
    margin = 72  # 1 inch
    line_height = 14
    max_chars = 95
    y = height - margin

    for paragraph in text.splitlines():
        paragraph = paragraph.rstrip()
        if not paragraph:
            y -= line_height
            if y <= margin:
                pdf_canvas.showPage()
                pdf_canvas.setFont("Helvetica", 11)
                y = height - margin
            continue

        for line in textwrap.wrap(paragraph, width=max_chars):
            pdf_canvas.drawString(margin, y, line)
            y -= line_height
            if y <= margin:
                pdf_canvas.showPage()
                pdf_canvas.setFont("Helvetica", 11)
                y = height - margin

    pdf_canvas.save()
    logger.info(f"✅ Fallback text-based PDF written to {output_path}")
    return True


def _html_to_text(html_content: str) -> str:
    if not html_content:
        return ""
    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
    except Exception:
        text = re.sub(r"(?i)<br\s*/?>", "\n", html_content)
        text = re.sub(r"<[^>]+>", " ", text)

    cleaned_lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in cleaned_lines if line)


def _fallback_convert_markdown_text(markdown_content: str, output_path: str, title: str = "Document") -> bool:
    text = _strip_markdown_syntax(markdown_content)
    return _fallback_convert_text_to_pdf(text, output_path, title)


def _strip_markdown_syntax(markdown_content: str) -> str:
    text = markdown_content
    text = re.sub(r"!\[[^\]]*\]\([^\)]*\)", "", text)  # images
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)  # links
    text = re.sub(r"`{1,3}[^`]+`{1,3}", "", text, flags=re.DOTALL)  # inline code
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)  # headers
    text = re.sub(r"[*_]{1,2}([^*_]+)[*_]{1,2}", r"\1", text)  # emphasis
    text = re.sub(r">+\s*", "", text)  # blockquotes
    text = re.sub(r"-{3,}", "", text)  # horizontal rules
    return text


def _convert_markdown_via_gotenberg(markdown_content, output_path, title="Document"):
    """
    Fallback function to convert Markdown using Gotenberg's native Markdown endpoint.
    This is used when the Python markdown library is not available.
    
    Args:
        markdown_content (str): Markdown content or file path
        output_path (str): Path where the PDF will be saved
        title (str): Title for the document
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    GOTENBERG_BASE_URL = os.getenv("PDF_CONVERT_URL", "http://localhost:7777")
    # Try Gotenberg v7+ Markdown endpoint
    api_url = f"{GOTENBERG_BASE_URL}/forms/chromium/convert/markdown"
    
    try:
        logger.info(f"Attempting direct Markdown conversion via Gotenberg at {output_path}")
        
        # Check if markdown_content is a file path or actual content
        if os.path.exists(markdown_content) and markdown_content.endswith('.md'):
            with open(markdown_content, 'r', encoding='utf-8') as f:
                md_content = f.read()
        else:
            md_content = markdown_content
        
        # Create the index.html template that Gotenberg expects for Markdown conversion
        index_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Arial, sans-serif;
            line-height: 1.6;
            max-width: 900px;
            margin: 0 auto;
            padding: 20px;
        }}
    </style>
</head>
<body>
    {{{{ toHTML .DirPath "content.md" }}}}
</body>
</html>"""
        
        # Prepare files for multipart upload
        files = [
            ('files', ('index.html', index_html.encode('utf-8'), 'text/html')),
            ('files', ('content.md', md_content.encode('utf-8'), 'text/markdown'))
        ]
        
        # Gotenberg parameters
        payload = {
            'paperWidth': '8.27',
            'paperHeight': '11.7',
            'marginTop': '0.39',
            'marginBottom': '0.39',
            'marginLeft': '0.39',
            'marginRight': '0.39',
            'printBackground': 'true',
            'landscape': 'false'
        }
        
        response = requests.post(
            api_url,
            files=files,
            data=payload,
            timeout=60
        )
        
        if response.status_code == 200:
            os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
            with open(output_path, 'wb') as out_file:
                out_file.write(response.content)
            logger.info(f"✅ Markdown to PDF conversion via Gotenberg successful: {output_path}")
            return True
        else:
            logger.error(f"❌ Gotenberg Markdown conversion failed with status {response.status_code}: {response.text}")
            return False
            
    except Exception as e:
        exc = traceback.format_exc()
        logger.error(f"❌ Exception in Gotenberg Markdown conversion: {e}\n{exc}")
        return False


def batch_convert_documents(file_list, output_dir):  
    """  
    Convert multiple documents to PDF using appropriate conversion function.  
      
    Args:  
        file_list (list): List of file paths to convert  
        output_dir (str): Directory to save converted PDFs  
          
    Returns:  
        dict: Conversion results with success/failure counts  
    """  
    results = {"success": 0, "failed": 0, "errors": []}  
      
    for file_path in file_list:  
        try:  
            # Determine output path  
            base_name = os.path.splitext(os.path.basename(file_path))[0]  
            output_path = os.path.join(output_dir, f"{base_name}.pdf")  
              
            # Choose appropriate conversion function  
            if file_path.lower().endswith(('.html', '.htm')):  
                success = convert_html_to_pdf(file_path, output_path)  
            elif file_path.lower().endswith(('.docx', '.doc', '.odt', '.rtf')):  
                success = convert_doc_to_pdf(file_path, output_path)  
            elif file_path.lower().endswith(('.md', '.markdown', '.mdown', '.mkd', '.mdwn')):  
                success = convert_markdown_file_to_pdf(file_path, output_path)  
            else:  
                results["errors"].append(f"Unsupported file type: {file_path}")  
                results["failed"] += 1  
                continue  
              
            if success:  
                results["success"] += 1  
            else:  
                results["failed"] += 1  
                results["errors"].append(f"Conversion failed: {file_path}")  
                  
        except Exception as e:  
            results["failed"] += 1  
            results["errors"].append(f"Error processing {file_path}: {str(e)}")  
      
    return results


def convert_markdown_to_pdf(markdown_content, output_path, title="Document", css_style=None):
    """
    Convert Markdown (file path or string content) to PDF.
    
    This is a convenience function that automatically detects whether the input
    is a file path or string content and calls the appropriate converter.
    
    Args:
        markdown_content (str): Markdown content as string or path to markdown file
        output_path (str): Path where the PDF will be saved
        title (str): Title for the document (used only for string content)
        css_style (str, optional): Custom CSS styles (used only for string content)
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    # Check if it's a file path
    if os.path.exists(markdown_content) and markdown_content.lower().endswith(('.md', '.markdown', '.mdown', '.mkd', '.mdwn')):
        return convert_markdown_file_to_pdf(markdown_content, output_path)
    else:
        # Treat as string content
        return convert_markdown_string_to_pdf(markdown_content, output_path, title, css_style)


def batch_convert_strings_to_pdf(content_list, output_dir):
    """
    Convert multiple content strings (HTML or Markdown) to PDF files.
    
    Args:
        content_list (list): List of dictionaries with keys:
            - 'content': The string content (HTML or Markdown)
            - 'type': Either 'html' or 'markdown'
            - 'filename': Name for the output PDF (without .pdf extension)
            - 'title' (optional): Document title
            - 'css_style' (optional): Custom CSS for markdown conversion
        output_dir (str): Directory to save converted PDFs
        
    Returns:
        dict: Conversion results with success/failure counts
        
    Example:
        content_list = [
            {'content': '<h1>Test</h1>', 'type': 'html', 'filename': 'test1'},
            {'content': '# Test', 'type': 'markdown', 'filename': 'test2', 'title': 'My Doc'}
        ]
        results = batch_convert_strings_to_pdf(content_list, './output')
    """
    results = {"success": 0, "failed": 0, "errors": []}
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    for item in content_list:
        try:
            # Validate required fields
            if 'content' not in item or 'type' not in item or 'filename' not in item:
                results["errors"].append(f"Missing required fields in item: {item}")
                results["failed"] += 1
                continue
            
            content = item['content']
            content_type = item['type'].lower()
            filename = item['filename']
            title = item.get('title', 'Document')
            css_style = item.get('css_style', None)
            
            # Ensure filename doesn't have .pdf extension (we'll add it)
            if filename.endswith('.pdf'):
                filename = filename[:-4]
            
            output_path = os.path.join(output_dir, f"{filename}.pdf")
            
            # Choose appropriate conversion function
            if content_type == 'html':
                success = convert_html_string_to_pdf(content, output_path, title)
            elif content_type == 'markdown':
                success = convert_markdown_string_to_pdf(content, output_path, title, css_style)
            else:
                results["errors"].append(f"Unsupported content type '{content_type}' for {filename}")
                results["failed"] += 1
                continue
            
            if success:
                results["success"] += 1
                logger.info(f"Successfully converted {filename} ({content_type}) to PDF")
            else:
                results["failed"] += 1
                results["errors"].append(f"Conversion failed for {filename}")
                
        except Exception as e:
            results["failed"] += 1
            error_msg = f"Error processing {item.get('filename', 'unknown')}: {str(e)}"
            results["errors"].append(error_msg)
            logger.error(error_msg)
    
    return results 